import streamlit as st
import os
import pickle
from langchain_groq import ChatGroq
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.schema import Document
from dotenv import load_dotenv
import time
import PyPDF2
from bs4 import BeautifulSoup
from pptx import Presentation
import pandas as pd
import openpyxl
from sentence_transformers import SentenceTransformer

# Load environment variables from .env file
load_dotenv()

# Set the title of the Streamlit app
st.title("Chatgroq With SentenceTransformer Embeddings Demo")

# Initialize the LLM
groq_api_key = os.getenv('GROQ_API_KEY')
if not groq_api_key:
    st.error("GROQ_API_KEY not found in environment variables. Please check your .env file.")
else:
    llm = ChatGroq(groq_api_key=groq_api_key)

# Initialize the SentenceTransformer model
model = SentenceTransformer('all-MiniLM-L6-v2')

# Create a HuggingFaceEmbeddings instance
embeddings = HuggingFaceEmbeddings(model_name='all-MiniLM-L6-v2')

# Function to load various types of documents
def load_documents(directory):
    documents = []
    for filename in os.listdir(directory):
        filepath = os.path.join(directory, filename)
        try:
            if filename.endswith(".pdf"):
                documents.extend(load_pdf_file(filepath))
            elif filename.endswith(".txt"):
                documents.extend(load_text_file(filepath))
            elif filename.endswith(".csv"):
                documents.extend(load_csv_file(filepath))
            elif filename.endswith(".xlsx") or filename.endswith(".xls"):
                documents.extend(load_excel_file(filepath))
            elif filename.endswith(".pptx"):
                documents.extend(load_pptx_file(filepath))
            elif filename.endswith(".html"):
                documents.extend(load_html_file(filepath))
        except Exception as e:
            st.warning(f"Error processing {filename}: {str(e)}")
    return documents

# Function to load PDF files
def load_pdf_file(filepath):
    documents = []
    with open(filepath, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page_number in range(len(reader.pages)):
            page = reader.pages[page_number]
            text = page.extract_text()
            if text:
                documents.append(Document(page_content=text, metadata={"source": filepath, "page_number": page_number + 1}))
    return documents

# Function to load text files
def load_text_file(filepath):
    with open(filepath, "r", encoding="utf-8") as file:
        text = file.read()
        return [Document(page_content=text, metadata={"source": filepath})]

# Function to load CSV files
def load_csv_file(filepath):
    df = pd.read_csv(filepath)
    text = df.to_string()
    return [Document(page_content=text, metadata={"source": filepath})]

# Function to load Excel files
def load_excel_file(filepath):
    df = pd.read_excel(filepath, engine='openpyxl')
    text = df.to_string()
    return [Document(page_content=text, metadata={"source": filepath})]

# Function to load PPTX files
def load_pptx_file(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return [Document(page_content=text, metadata={"source": filepath})]

# Function to load HTML files
def load_html_file(filepath):
    with open(filepath, "r", encoding="utf-8") as file:
        soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text(separator="\n")
        return [Document(page_content=text, metadata={"source": filepath})]

# Function to initialize or load vector embeddings
def initialize_or_load_vector_embeddings():
    faiss_index_path = "faiss_index"
    metadata_path = "document_metadata.pkl"
    
    if os.path.exists(faiss_index_path) and os.path.exists(metadata_path):
        st.info("Loading existing embeddings and metadata...")
        faiss_index = FAISS.load_local(faiss_index_path, embeddings)
        with open(metadata_path, 'rb') as f:
            st.session_state.final_documents = pickle.load(f)
        st.success(f"Loaded {len(st.session_state.final_documents)} documents from saved data.")
        return faiss_index
    
    st.info("No existing embeddings found. Initializing new embeddings...")
    return initialize_vector_embeddings()

# Function to initialize vector embeddings
def initialize_vector_embeddings():
    try:
        data_directory = r"C:\Users\harsh\Downloads\chatbot\data"
        log_directory = r"C:\Users\harsh\Downloads\chatbot\log"

        if not os.path.exists(data_directory) or not os.path.exists(log_directory):
            st.error(f"One or both directories do not exist: {data_directory}, {log_directory}")
            return None

        data_documents = load_documents(data_directory)
        log_documents = load_documents(log_directory)
        all_documents = data_documents + log_documents

        if not all_documents:
            st.error("No documents loaded. Please check your document loading functions and directory contents.")
            return None

        texts = [doc.page_content for doc in all_documents]
        
        faiss_index = FAISS.from_documents(all_documents, embeddings)
        
        # Save the FAISS index
        faiss_index.save_local("faiss_index")
        
        # Save document metadata
        with open("document_metadata.pkl", 'wb') as f:
            pickle.dump(all_documents, f)
        
        st.session_state.final_documents = all_documents
        st.success(f"Number of documents processed and saved: {len(st.session_state.final_documents)}")
        
        return faiss_index

    except Exception as e:
        st.error(f"Error during vector embedding: {str(e)}")
        return None

# Main app logic
if 'faiss_index' not in st.session_state:
    st.session_state.faiss_index = None

if st.button("Initialize or Load Documents Embedding"):
    st.session_state.faiss_index = initialize_or_load_vector_embeddings()

prompt1 = st.text_input("Enter Your Question From Documents")

if prompt1 and st.session_state.faiss_index:
    try:
        prompt = ChatPromptTemplate.from_template(
            """
            Answer the questions based on the provided context only.
            Please provide the most accurate response based on the question.
            <context>
            {context}
            </context>
            Question: {input}
            """
        )

        document_chain = create_stuff_documents_chain(llm, prompt)
        retrieval_chain = create_retrieval_chain(st.session_state.faiss_index.as_retriever(), document_chain)

        start = time.process_time()
        response = retrieval_chain.invoke({'input': prompt1})
        st.write("Response time:", time.process_time() - start)
        st.write(response['answer'])

        with st.expander("Document Similarity Search"):
            for i, doc in enumerate(response["context"]):
                st.write(f"Content from {doc.metadata.get('source', 'Unknown')}:")
                st.write(doc.page_content[:200] + "...")  # Show first 200 characters
                st.write("--------------------------------")

    except Exception as e:
        st.error(f"Error during retrieval or processing: {str(e)}")

if "final_documents" in st.session_state:
    st.write(f"Number of processed documents: {len(st.session_state.final_documents)}")

# Add a button to force reprocessing
if st.button("Force Reprocess All Documents"):
    if os.path.exists("faiss_index"):
        import shutil
        shutil.rmtree("faiss_index")
    if os.path.exists("document_metadata.pkl"):
        os.remove("document_metadata.pkl")
    st.session_state.faiss_index = initialize_vector_embeddings()