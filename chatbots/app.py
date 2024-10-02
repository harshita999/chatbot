import streamlit as st
import os
from langchain_groq import ChatGroq
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, CSVLoader, UnstructuredExcelLoader
from langchain.schema import Document
from dotenv import load_dotenv
import time
import PyPDF2
from bs4 import BeautifulSoup
import pandas as pd
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

load_dotenv()

groq_api_key = os.getenv('GROQ_API_KEY')

st.title("Chatgroq With Llama3 Demo")

llm = ChatGroq(groq_api_key=groq_api_key, model_name="Llama3-8b-8192")

prompt = ChatPromptTemplate.from_template(
    """
    Answer the questions based on the provided context only.
    Please provide the most accurate response based on the question.
    <context>
    {context}
    <context>
    Questions:{input}
    """
)

def load_html_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".html"):
            filepath = os.path.join(directory, filename)
            with open(filepath, "r", encoding="utf-8") as file:
                soup = BeautifulSoup(file, "html.parser")
                text = soup.get_text(separator="\n")
                documents.append(Document(page_content=text, metadata={"source": filepath}))
    return documents

def load_text_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".txt") and not filename.startswith("~$"):  # Skip temporary/system files
            filepath = os.path.join(directory, filename)
            try:
                loader = TextLoader(filepath)
                text_documents = loader.load()
                documents.extend(text_documents)
            except Exception as e:
                st.error(f"Error loading text file {filepath}: {e}")
    return documents

def load_csv_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".csv") and not filename.startswith("~$"):  # Skip temporary/system files
            filepath = os.path.join(directory, filename)
            try:
                loader = CSVLoader(filepath)
                csv_documents = loader.load()
                documents.extend(csv_documents)
            except Exception as e:
                st.error(f"Error loading CSV file {filepath}: {e}")
    return documents

def load_excel_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".xlsx") and not filename.startswith("~$"):  # Skip temporary/system files
            filepath = os.path.join(directory, filename)
            try:
                loader = UnstructuredExcelLoader(filepath)
                excel_documents = loader.load()
                documents.extend(excel_documents)
            except Exception as e:
                st.error(f"Error loading Excel file {filepath}: {e}")
    return documents

def load_pptx_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".pptx") and not filename.startswith("~$"):  # Skip temporary/system files
            filepath = os.path.join(directory, filename)
            try:
                presentation = Presentation(filepath)  # Use python-pptx to load the file
                text_content = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                documents.append(Document(page_content=text_content, metadata={"source": filepath}))
            except Exception as e:
                st.error(f"Error loading PPTX file {filepath}: {e}")
    return documents

def load_pdf_files(directory):
    documents = []
    for filename in os.listdir(directory):
        if filename.endswith(".pdf") and not filename.startswith("~$"):  # Skip temporary/system files
            filepath = os.path.join(directory, filename)
            try:
                with open(filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    num_pages = len(reader.pages)
                    for page_number in range(num_pages):
                        page = reader.pages[page_number]
                        text = page.extract_text()
                        if text:  # Ensure the page has content
                            documents.append(Document(page_content=text, metadata={"source": filepath, "page_number": page_number + 1}))
            except Exception as e:
                st.error(f"Error reading PDF {filepath}: {e}")
    return documents

def initialize_vector_embeddings():
    try:
        st.session_state.embeddings = OpenAIEmbeddings()

        # Load documents from various formats
        st.session_state.pdf_docs = load_pdf_files("./data")
        st.session_state.txt_docs = load_text_files("./data")
        st.session_state.csv_docs = load_csv_files("./data")
        st.session_state.excel_docs = load_excel_files("./data")
        st.session_state.pptx_docs = load_pptx_files("./data")
        st.session_state.html_docs = load_html_files("./log")

        # Debug: Print document counts
        st.write(f"PDF Documents: {len(st.session_state.pdf_docs)}")
        st.write(f"Text Documents: {len(st.session_state.txt_docs)}")
        st.write(f"CSV Documents: {len(st.session_state.csv_docs)}")
        st.write(f"Excel Documents: {len(st.session_state.excel_docs)}")
        st.write(f"PPTX Documents: {len(st.session_state.pptx_docs)}")
        st.write(f"HTML Documents: {len(st.session_state.html_docs)}")

        # Text splitter
        st.session_state.text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=2000)

        all_docs = (st.session_state.pdf_docs +
                    st.session_state.txt_docs +
                    st.session_state.csv_docs +
                    st.session_state.excel_docs +
                    st.session_state.pptx_docs +
                    st.session_state.html_docs)

        documents_with_page_numbers = []

        for doc in all_docs:
            chunks = st.session_state.text_splitter.split_text(doc.page_content)
            for chunk in chunks:
                chunk_document = Document(
                    page_content=chunk,
                    metadata={"source": doc.metadata["source"]}
                )
                documents_with_page_numbers.append(chunk_document)

        st.session_state.final_documents = documents_with_page_numbers

        st.write(f"Number of documents processed: {len(st.session_state.final_documents)}")
        if st.session_state.final_documents:
            sample_doc = st.session_state.final_documents[0]
            st.write("Sample document content:", sample_doc.page_content[:200])
            st.write("Source:", sample_doc.metadata["source"])
            if 'page_number' in sample_doc.metadata:
                st.write("Page number:", sample_doc.metadata["page_number"])

        st.session_state.vectors = FAISS.from_documents(st.session_state.final_documents, st.session_state.embeddings)

        vector_store_path = "./vector_store/faiss_index"
        st.session_state.vectors.save(vector_store_path)
        st.write("Vector Store DB Is Ready and Saved to Disk")
    except Exception as e:
        st.error(f"Error during vector embedding: {e}")

def load_vector_embeddings():
    vector_store_path = "./vector_store/faiss_index"

    if os.path.exists(vector_store_path):
        try:
            st.session_state.vectors = FAISS.load(vector_store_path)
            st.write("Vector Store DB Loaded from Disk")
        except Exception as e:
            st.error(f"Error loading vector embeddings from disk: {e}")
    else:
        st.error("Vector embeddings are not initialized and no saved embeddings found on disk. Please initialize the embeddings.")

if "vectors" not in st.session_state:
    load_vector_embeddings()

if st.button("Initialize Documents Embedding"):
    initialize_vector_embeddings()

prompt1 = st.text_input("Enter Your Question From Documents")

if prompt1:
    if "vectors" in st.session_state and st.session_state.vectors:
        try:
            document_chain = create_stuff_documents_chain(llm, prompt)
            retriever = st.session_state.vectors.as_retriever()
            retrieval_chain = create_retrieval_chain(retriever, document_chain)
            start = time.process_time()
            response = retrieval_chain.invoke({'input': prompt1})
            st.write("Response time:", time.process_time() - start)
            st.write(response['answer'])

            with st.expander("Document Similarity Search"):
                for i, doc in enumerate(response["context"]):
                    st.write(f"Content from {doc.metadata.get('source', 'Unknown')}:")
                    st.write(doc.page_content)
                    st.write("--------------------------------")
        except Exception as e:
            st.error(f"Error during retrieval or processing: {e}")
    else:
        st.error("Vector embeddings are not initialized. Please click the 'Initialize Documents Embedding' button first.")
